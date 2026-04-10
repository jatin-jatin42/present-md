"""PPTX Builder — assembles the final presentation from the visual plan."""

from __future__ import annotations

import logging
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from present_md.visual.decision import PresentationPlan, SlidePlan, VisualElement
from present_md.visual.charts import add_chart_to_slide, add_table_to_slide
from present_md.renderer.template import ThemeConfig
from present_md.layout.grid import GridSystem, Rect

logger = logging.getLogger(__name__)


class PresentationBuilder:
    """Builds the final PPTX file from a PresentationPlan."""

    def __init__(self, theme: ThemeConfig, grid: GridSystem, template_path: str):
        self.theme = theme
        self.grid = grid
        self.prs = Presentation(template_path)
        # Remove any existing slides from the template
        self._clear_existing_slides()

    def _clear_existing_slides(self):
        """Remove existing slides from the template, keeping only masters/layouts."""
        for i in range(len(self.prs.slides) - 1, -1, -1):
            slide_id = self.prs.slides._sldIdLst[i].id
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]

    def build(self, plan: PresentationPlan):
        """Build all slides from the plan."""
        for slide_plan in plan.slides:
            try:
                self._build_slide(slide_plan)
            except Exception as e:
                logger.error(f"Error building slide {slide_plan.slide_number}: {e}")
                # Fallback: create a simple text slide
                self._build_fallback_slide(slide_plan)

    def save(self, output_path: str):
        """Save the presentation to disk."""
        self.prs.save(output_path)

    def _get_blank_layout(self):
        """Get a blank slide layout (or the first available one)."""
        # Try to find a blank layout
        for layout in self.prs.slide_layouts:
            if "blank" in layout.name.lower():
                return layout
        # Fallback to last layout (usually blank-ish)
        return self.prs.slide_layouts[-1]

    def _get_title_layout(self):
        """Get a title slide layout."""
        for layout in self.prs.slide_layouts:
            name = layout.name.lower()
            if "title" in name and ("slide" in name or "layout" in name):
                return layout
        # Fallback to first layout
        return self.prs.slide_layouts[0]

    def _build_slide(self, plan: SlidePlan):
        """Build a single slide based on its plan."""
        if plan.slide_type == "title":
            self._build_title_slide(plan)
        elif plan.slide_type == "agenda":
            self._build_agenda_slide(plan)
        elif plan.slide_type == "executive_summary":
            self._build_content_slide(plan)
        elif plan.slide_type == "conclusion":
            self._build_content_slide(plan)
        else:
            self._build_content_slide(plan)

    def _build_title_slide(self, plan: SlidePlan):
        """Build the title slide."""
        layout = self._get_title_layout()
        slide = self.prs.slides.add_slide(layout)

        # Clear default placeholders
        for shape in list(slide.placeholders):
            sp = shape._element
            sp.getparent().remove(sp)

        # Title
        title_rect = self.grid.title_rect
        title_rect = Rect(
            title_rect.left, Inches(2.5),
            title_rect.width, Inches(1.5),
        )
        self._add_text_box(
            slide, title_rect,
            plan.title,
            font_size=Pt(36),
            font_bold=True,
            font_color=self.theme.colors[0] if self.theme.colors else RGBColor(0, 0, 0),
            font_name=self.theme.title_font,
            alignment=PP_ALIGN.CENTER,
        )

        # Subtitle / key message
        if plan.key_message:
            sub_rect = Rect(
                title_rect.left, title_rect.bottom + Inches(0.3),
                title_rect.width, Inches(0.8),
            )
            self._add_text_box(
                slide, sub_rect,
                plan.key_message,
                font_size=Pt(18),
                font_color=self.theme.colors[2] if len(self.theme.colors) > 2 else RGBColor(0x66, 0x66, 0x66),
                font_name=self.theme.body_font,
                alignment=PP_ALIGN.CENTER,
            )

    def _build_agenda_slide(self, plan: SlidePlan):
        """Build the agenda/TOC slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Clear placeholders
        for shape in list(slide.placeholders):
            sp = shape._element
            sp.getparent().remove(sp)

        # Title
        self._add_slide_title(slide, plan.title)

        # Agenda items
        items = plan.elements[0].content.get("items", plan.key_message.split(",")) if plan.elements else []
        if not items:
            items = [plan.key_message] if plan.key_message else ["Overview"]

        content_rect = self.grid.content_rect
        self._add_bullet_list(slide, content_rect, items)

    def _build_content_slide(self, plan: SlidePlan):
        """Build a content slide based on its visual type."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Clear placeholders
        for shape in list(slide.placeholders):
            sp = shape._element
            sp.getparent().remove(sp)

        # Always add title
        self._add_slide_title(slide, plan.title)

        # Add key message subtitle if present
        if plan.key_message and plan.visual_type not in ("key_message_box",):
            self._add_key_message(slide, plan.key_message)

        # Render based on visual type
        visual_type = plan.visual_type

        if visual_type == "stat_callout":
            self._render_stat_callout(slide, plan)
        elif visual_type == "icon_grid":
            self._render_icon_grid(slide, plan)
        elif visual_type == "process_flow":
            self._render_process_flow(slide, plan)
        elif visual_type in ("chart_bar", "chart_pie", "chart_line", "chart_area"):
            self._render_chart(slide, plan)
        elif visual_type == "table":
            self._render_table(slide, plan)
        elif visual_type == "comparison":
            self._render_comparison(slide, plan)
        elif visual_type == "timeline":
            self._render_timeline(slide, plan)
        else:
            # Default: bullet list
            self._render_bullet_list(slide, plan)

    def _build_fallback_slide(self, plan: SlidePlan):
        """Build a simple fallback slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())
        for shape in list(slide.placeholders):
            sp = shape._element
            sp.getparent().remove(sp)
        self._add_slide_title(slide, plan.title)
        if plan.key_message:
            self._add_key_message(slide, plan.key_message)

    # ── Helper methods ─────────────────────────────────────────────────

    def _add_slide_title(self, slide, text: str):
        """Add a title to a slide."""
        rect = self.grid.title_rect
        self._add_text_box(
            slide, rect, text,
            font_size=Pt(24),
            font_bold=True,
            font_color=self.theme.colors[0] if self.theme.colors else RGBColor(0, 0, 0),
            font_name=self.theme.title_font,
            alignment=PP_ALIGN.LEFT,
        )
        # Add a thin accent line below title
        accent = self.theme.accent_colors[0] if self.theme.accent_colors else RGBColor(0x27, 0xAE, 0x60)
        line_shape = slide.shapes.add_shape(
            1,  # Rectangle
            rect.left, rect.bottom - Inches(0.05),
            Inches(1.5), Inches(0.04),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = accent
        line_shape.line.fill.background()

    def _add_key_message(self, slide, text: str):
        """Add a key message subtitle."""
        rect = self.grid.subtitle_rect
        self._add_text_box(
            slide, rect, text,
            font_size=Pt(12),
            font_italic=True,
            font_color=self.theme.colors[2] if len(self.theme.colors) > 2 else RGBColor(0x66, 0x66, 0x66),
            font_name=self.theme.body_font,
            alignment=PP_ALIGN.LEFT,
        )

    def _add_text_box(
        self, slide, rect: Rect, text: str,
        font_size=Pt(14), font_bold=False, font_italic=False,
        font_color=None, font_name=None,
        alignment=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP,
        fill_color=None,
    ):
        """Add a text box to a slide with consistent formatting."""
        txBox = slide.shapes.add_textbox(
            rect.left, rect.top, rect.width, rect.height
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Fill
        if fill_color:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = fill_color
        else:
            txBox.fill.background()

        # No line
        txBox.line.fill.background()

        # Text
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = font_bold
        p.font.italic = font_italic
        if font_color:
            p.font.color.rgb = font_color
        if font_name:
            p.font.name = font_name
        p.alignment = alignment

        return txBox

    def _add_bullet_list(self, slide, rect: Rect, items: list[str]):
        """Add a bulleted list to a slide."""
        txBox = slide.shapes.add_textbox(
            rect.left, rect.top, rect.width, rect.height
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        txBox.fill.background()
        txBox.line.fill.background()

        for i, item in enumerate(items[:8]):  # Max 8 items
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(self.theme.body_font_size)
            p.font.name = self.theme.body_font
            p.font.color.rgb = self.theme.colors[0] if self.theme.colors else RGBColor(0, 0, 0)
            p.level = 0
            p.space_after = Pt(8)
            # Add bullet
            p.bullet = True

        return txBox

    def _add_shape_box(
        self, slide, rect: Rect,
        fill_color: RGBColor = None,
        corner_radius: int = Inches(0.1),
    ):
        """Add a rounded rectangle container shape."""
        shape = slide.shapes.add_shape(
            5,  # Rounded rectangle
            rect.left, rect.top, rect.width, rect.height,
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        shape.line.fill.background()
        return shape

    # ── Visual Renderers ─────────────────────────────────────────────

    def _render_stat_callout(self, slide, plan: SlidePlan):
        """Render stat callout boxes."""
        elements = plan.elements
        if not elements:
            return self._render_bullet_list(slide, plan)

        content = elements[0].content
        stats = content.get("stats", content.get("items", []))
        if not stats:
            return self._render_bullet_list(slide, plan)

        n = min(len(stats), 4)
        rects = self.grid.stat_callout_layout(n)

        for i in range(n):
            rect = rects[i]
            color = self.theme.accent_colors[i % len(self.theme.accent_colors)] if self.theme.accent_colors else RGBColor(0x27, 0xAE, 0x60)

            # Container
            self._add_shape_box(slide, rect, fill_color=color)

            # Stat value
            if isinstance(stats[i], dict):
                value = str(stats[i].get("value", ""))
                label = stats[i].get("label", "")
            else:
                value = str(stats[i])
                label = ""

            val_rect = Rect(rect.left, rect.top + Inches(0.3), rect.width, Inches(1.0))
            self._add_text_box(
                slide, val_rect, value,
                font_size=Pt(32), font_bold=True,
                font_color=RGBColor(0xFF, 0xFF, 0xFF),
                font_name=self.theme.title_font,
                alignment=PP_ALIGN.CENTER,
            )

            if label:
                lbl_rect = Rect(rect.left, val_rect.bottom, rect.width, Inches(0.6))
                self._add_text_box(
                    slide, lbl_rect, label,
                    font_size=Pt(11),
                    font_color=RGBColor(0xFF, 0xFF, 0xFF),
                    font_name=self.theme.body_font,
                    alignment=PP_ALIGN.CENTER,
                )

    def _render_icon_grid(self, slide, plan: SlidePlan):
        """Render icon + label grid."""
        elements = plan.elements
        if not elements:
            return self._render_bullet_list(slide, plan)

        content = elements[0].content
        items = content.get("items", [])
        if not items:
            return self._render_bullet_list(slide, plan)

        n = min(len(items), 6)
        rects = self.grid.icon_grid_layout(n)

        for i in range(n):
            rect = rects[i]
            inner = rect.shrink(Inches(0.1))

            # Light background box
            bg_color = self.theme.colors[3] if len(self.theme.colors) > 3 else RGBColor(0xEC, 0xF0, 0xF1)
            self._add_shape_box(slide, inner, fill_color=bg_color)

            # Icon placeholder (colored circle with number)
            accent = self.theme.accent_colors[i % len(self.theme.accent_colors)] if self.theme.accent_colors else RGBColor(0x27, 0xAE, 0x60)
            icon_size = Inches(0.6)
            icon_rect = Rect(
                inner.left + (inner.width - icon_size) // 2,
                inner.top + Inches(0.3),
                icon_size, icon_size,
            )
            icon_shape = slide.shapes.add_shape(
                9,  # Oval
                icon_rect.left, icon_rect.top,
                icon_rect.width, icon_rect.height,
            )
            icon_shape.fill.solid()
            icon_shape.fill.fore_color.rgb = accent
            icon_shape.line.fill.background()

            # Number in icon
            if icon_shape.has_text_frame:
                tf = icon_shape.text_frame
                p = tf.paragraphs[0]
                p.text = str(i + 1)
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER

            # Label text
            item_text = items[i] if isinstance(items[i], str) else str(items[i])
            label_rect = Rect(
                inner.left + Inches(0.1),
                icon_rect.bottom + Inches(0.2),
                inner.width - Inches(0.2),
                inner.height - Inches(1.3),
            )
            self._add_text_box(
                slide, label_rect, item_text,
                font_size=Pt(10),
                font_color=self.theme.colors[0] if self.theme.colors else RGBColor(0, 0, 0),
                font_name=self.theme.body_font,
                alignment=PP_ALIGN.CENTER,
            )

    def _render_process_flow(self, slide, plan: SlidePlan):
        """Render a process flow (steps with arrows)."""
        elements = plan.elements
        if not elements:
            return self._render_bullet_list(slide, plan)

        content = elements[0].content
        steps = content.get("steps", content.get("items", []))
        if not steps:
            return self._render_bullet_list(slide, plan)

        n = min(len(steps), 6)
        rects = self.grid.process_flow_layout(n)

        for i in range(n):
            rect = rects[i]
            accent = self.theme.accent_colors[i % len(self.theme.accent_colors)] if self.theme.accent_colors else RGBColor(0x27, 0xAE, 0x60)

            # Step box
            self._add_shape_box(slide, rect, fill_color=accent)

            # Step number
            num_rect = Rect(rect.left, rect.top + Inches(0.1), rect.width, Inches(0.4))
            self._add_text_box(
                slide, num_rect, str(i + 1),
                font_size=Pt(20), font_bold=True,
                font_color=RGBColor(0xFF, 0xFF, 0xFF),
                font_name=self.theme.title_font,
                alignment=PP_ALIGN.CENTER,
            )

            # Step text
            step_text = steps[i] if isinstance(steps[i], str) else str(steps[i])
            text_rect = Rect(
                rect.left + Inches(0.1), num_rect.bottom,
                rect.width - Inches(0.2), rect.height - Inches(0.6),
            )
            self._add_text_box(
                slide, text_rect, step_text,
                font_size=Pt(9),
                font_color=RGBColor(0xFF, 0xFF, 0xFF),
                font_name=self.theme.body_font,
                alignment=PP_ALIGN.CENTER,
            )

            # Arrow between steps (except last)
            if i < n - 1:
                arrow_left = rect.right + Inches(0.02)
                arrow_top = rect.top + rect.height // 2 - Inches(0.1)
                self._add_text_box(
                    slide,
                    Rect(arrow_left, arrow_top, Inches(0.15), Inches(0.2)),
                    "→",
                    font_size=Pt(16), font_bold=True,
                    font_color=self.theme.colors[0] if self.theme.colors else RGBColor(0, 0, 0),
                    alignment=PP_ALIGN.CENTER,
                )

    def _render_chart(self, slide, plan: SlidePlan):
        """Render a chart."""
        elements = plan.elements
        if not elements:
            return self._render_bullet_list(slide, plan)

        content = elements[0].content
        headers = content.get("headers", [])
        rows = content.get("rows", [])

        if not headers or not rows:
            return self._render_bullet_list(slide, plan)

        chart_rect = self.grid.chart_layout()
        add_chart_to_slide(
            slide,
            plan.visual_type,
            headers, rows,
            title=content.get("title", plan.title),
            left=chart_rect.left, top=chart_rect.top,
            width=chart_rect.width, height=chart_rect.height,
            colors=self.theme.accent_colors,
        )

    def _render_table(self, slide, plan: SlidePlan):
        """Render a data table."""
        elements = plan.elements
        if not elements:
            return self._render_bullet_list(slide, plan)

        content = elements[0].content
        headers = content.get("headers", [])
        rows = content.get("rows", [])

        if not headers or not rows:
            return self._render_bullet_list(slide, plan)

        table_rect = self.grid.table_layout()
        add_table_to_slide(
            slide, headers, rows,
            title=content.get("title", plan.title),
            left=table_rect.left, top=table_rect.top,
            width=table_rect.width,
            colors=self.theme.accent_colors[:2] if self.theme.accent_colors else None,
            font_size=self.theme.body_font_size,
        )

    def _render_comparison(self, slide, plan: SlidePlan):
        """Render a comparison layout (two columns)."""
        left_col, right_col = self.grid.two_column_layout()

        elements = plan.elements
        if not elements or len(elements) < 2:
            return self._render_bullet_list(slide, plan)

        for i, (rect, elem) in enumerate(zip([left_col, right_col], elements[:2])):
            content = elem.content
            col_title = content.get("title", f"Option {i+1}")
            items = content.get("items", [])

            accent = self.theme.accent_colors[i % len(self.theme.accent_colors)] if self.theme.accent_colors else RGBColor(0x27, 0xAE, 0x60)

            # Column header
            header_rect = Rect(rect.left, rect.top, rect.width, Inches(0.5))
            self._add_shape_box(slide, header_rect, fill_color=accent)
            self._add_text_box(
                slide, header_rect, col_title,
                font_size=Pt(14), font_bold=True,
                font_color=RGBColor(0xFF, 0xFF, 0xFF),
                font_name=self.theme.title_font,
                alignment=PP_ALIGN.CENTER,
            )

            # Column items
            items_rect = Rect(
                rect.left, header_rect.bottom + Inches(0.1),
                rect.width, rect.height - Inches(0.6),
            )
            self._add_bullet_list(slide, items_rect, items[:5])

    def _render_timeline(self, slide, plan: SlidePlan):
        """Render a timeline (horizontal steps)."""
        # Reuse process flow with slight style adjustment
        self._render_process_flow(slide, plan)

    def _render_bullet_list(self, slide, plan: SlidePlan):
        """Render a simple bullet list (fallback)."""
        elements = plan.elements
        items = []

        if elements:
            content = elements[0].content
            items = content.get("items", [])

        if not items:
            # Try to extract from key_message
            if plan.key_message:
                items = [plan.key_message]

        if items:
            bullet_rect = self.grid.bullet_layout()
            self._add_bullet_list(slide, bullet_rect, items[:6])
