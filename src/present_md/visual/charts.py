"""Chart and table generation utilities for python-pptx."""

from __future__ import annotations

import logging
import re
from typing import Optional

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


def parse_numeric(value: str) -> Optional[float]:
    """Try to parse a string as a number, handling $, %, commas, etc."""
    if not value:
        return None
    cleaned = re.sub(r'[,$%\sA-Za-z]', '', value.strip())
    # Handle special chars like ‑ (non-breaking hyphen)
    cleaned = cleaned.replace('‑', '-').replace('–', '-')
    try:
        return float(cleaned)
    except ValueError:
        return None


def add_chart_to_slide(
    slide,
    chart_type: str,  # "bar", "pie", "line", "area"
    headers: list[str],
    rows: list[list[str]],
    title: str = "",
    left: int = None,
    top: int = None,
    width: int = None,
    height: int = None,
    colors: list[RGBColor] = None,
):
    """Add a native PPTX chart to a slide.

    Args:
        slide: pptx slide object
        chart_type: one of "bar", "pie", "line", "area"
        headers: column headers
        rows: data rows
        title: chart title
        left, top, width, height: position/size in EMUs
        colors: theme colors to apply
    """
    left = left or Inches(0.8)
    top = top or Inches(1.8)
    width = width or Inches(8.5)
    height = height or Inches(4.5)

    # Determine chart type enum
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "chart_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
        "chart_pie": XL_CHART_TYPE.PIE,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "chart_line": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Build chart data
    chart_data = CategoryChartData()

    if len(headers) < 2 or not rows:
        logger.warning("Not enough data for chart, skipping")
        return None

    # Categories from first column
    categories = [row[0] if row else "" for row in rows]
    chart_data.categories = categories

    # Series from remaining columns
    for col_idx in range(1, len(headers)):
        series_name = headers[col_idx]
        values = []
        for row in rows:
            if col_idx < len(row):
                val = parse_numeric(row[col_idx])
                values.append(val if val is not None else 0)
            else:
                values.append(0)
        chart_data.add_series(series_name, values)

    # Add chart to slide
    try:
        chart_frame = slide.shapes.add_chart(
            xl_chart_type, left, top, width, height, chart_data
        )
        chart = chart_frame.chart

        # Style the chart
        chart.has_legend = True
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # Apply colors if provided
        if colors and chart.series:
            for i, series in enumerate(chart.series):
                if i < len(colors):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = colors[i]

        return chart_frame
    except Exception as e:
        logger.error(f"Failed to create chart: {e}")
        return None


def add_table_to_slide(
    slide,
    headers: list[str],
    rows: list[list[str]],
    title: str = "",
    left: int = None,
    top: int = None,
    width: int = None,
    height: int = None,
    colors: list[RGBColor] = None,
    font_size: int = 10,
):
    """Add a formatted table to a slide.

    Args:
        slide: pptx slide object
        headers: column headers
        rows: data rows (max 8 rows for readability)
        left, top, width, height: position/size in EMUs
        colors: theme colors (header_bg, alt_row_bg)
        font_size: font size in points
    """
    left = left or Inches(0.8)
    top = top or Inches(1.8)
    width = width or Inches(8.5)

    # Limit rows for readability
    display_rows = rows[:8]
    total_rows = len(display_rows) + 1  # +1 for header
    total_cols = len(headers)

    if total_cols == 0:
        return None

    row_height = Inches(0.4)
    height = height or (row_height * total_rows)

    try:
        table_shape = slide.shapes.add_table(
            total_rows, total_cols, left, top, width, height
        )
        table = table_shape.table

        # Style header row
        header_color = colors[0] if colors else RGBColor(0x2C, 0x3E, 0x50)
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            # Format
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                paragraph.alignment = 1  # Center
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            # Vertical alignment — middle
            cell.vertical_anchor = 1  # Middle

        # Style data rows
        alt_color = colors[1] if colors and len(colors) > 1 else RGBColor(0xEC, 0xF0, 0xF1)
        for row_idx, row_data in enumerate(display_rows):
            for col_idx in range(total_cols):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = row_data[col_idx] if col_idx < len(row_data) else ""
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size)
                    paragraph.alignment = 1  # Center
                cell.vertical_anchor = 1  # Middle

                # Alternate row coloring
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_color

        return table_shape
    except Exception as e:
        logger.error(f"Failed to create table: {e}")
        return None
